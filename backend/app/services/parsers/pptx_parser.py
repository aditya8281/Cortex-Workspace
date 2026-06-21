"""PowerPoint presentation (.pptx) parser."""

from __future__ import annotations

import logging

from backend.app.services.parsers.base import BaseParser, ParsedDocument, ParsedSection

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    """Extract text from PowerPoint .pptx files using python-pptx."""

    def parse(self, file_path: str) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not installed — run: pip install python-pptx")
            return ParsedDocument(metadata={"error": "python-pptx not installed"})

        sections: list[ParsedSection] = []
        metadata: dict = {}

        try:
            presentation = Presentation(file_path)

            if presentation.core_properties.title:
                metadata["title"] = presentation.core_properties.title
            if presentation.core_properties.author:
                metadata["author"] = presentation.core_properties.author
            metadata["slide_count"] = len(presentation.slides)

            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text_parts: list[str] = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_text_parts.append(text)

                    if shape.has_table:
                        table_text = self._format_table(shape.table)
                        if table_text.strip():
                            slide_text_parts.append(table_text)

                # Add speaker notes if present
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    notes_text = notes_frame.text.strip()
                    if notes_text:
                        slide_text_parts.append(f"[Notes: {notes_text}]")

                if slide_text_parts:
                    content = "\n\n".join(slide_text_parts)
                    sections.append(
                        ParsedSection(
                            title=f"Slide {slide_num}",
                            content=content,
                            section_type="paragraph",
                        )
                    )

        except Exception as e:
            logger.warning("Failed to parse PPTX %s: %s", file_path, e)
            return ParsedDocument(metadata={"error": str(e)})

        total_chars = sum(len(s.content) for s in sections)
        return ParsedDocument(sections=sections, metadata=metadata, total_chars=total_chars)

    @staticmethod
    def _format_table(table) -> str:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        headers = rows[0]
        data_rows = rows[1:]
        header_line = " | ".join(headers)
        separator = " | ".join(["---"] * len(headers))
        lines = [header_line, separator]
        for row in data_rows:
            lines.append(" | ".join(row))
        return "\n".join(lines)
